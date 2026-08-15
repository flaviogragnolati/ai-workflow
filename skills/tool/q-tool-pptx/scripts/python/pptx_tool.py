#!/usr/bin/env python3
"""Python PPTX backend for q-tool-pptx.

Unified command-line adapter over python-pptx, Pillow, LibreOffice, and
Poppler for inspecting, extracting, restructuring, and validating PPTX/POTX
packages. Slide numbers are 1-based in every user-facing argument. Outputs are
written atomically and never replace an existing file or non-empty directory
unless --overwrite is passed explicitly.
"""

from __future__ import annotations

import argparse
import json
import os
import posixpath
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

EXIT_OK = 0
EXIT_ARGS = 2
EXIT_INPUT = 3
EXIT_MISSING_DEPENDENCY = 4
EXIT_UNSUPPORTED = 5
EXIT_OUTPUT = 6
EXIT_PARTIAL = 8

PPTX_EXTENSIONS = {".pptx", ".potx", ".ppsx"}
MAX_ARCHIVE_BYTES = 64 * 1024 * 1024
MAX_ENTRIES = 4096
MAX_EXPANDED_BYTES = 256 * 1024 * 1024
MAX_COMPRESSION_RATIO = 250
MAX_XML_BYTES = 16 * 1024 * 1024


class ToolError(Exception):
    def __init__(self, message: str, code: int = EXIT_OUTPUT) -> None:
        super().__init__(message)
        self.code = code


def load_pptx_module():
    try:
        import pptx  # noqa: PLC0415

        return pptx
    except ModuleNotFoundError as exc:
        raise ToolError(
            "python-pptx is not installed; run `pip install -e .` inside scripts/python",
            EXIT_MISSING_DEPENDENCY,
        ) from exc


def load_pillow_module():
    try:
        from PIL import Image, ImageDraw  # noqa: PLC0415

        return Image, ImageDraw
    except ModuleNotFoundError as exc:
        raise ToolError(
            "Pillow is not installed; run `pip install -e .` inside scripts/python",
            EXIT_MISSING_DEPENDENCY,
        ) from exc


def xml_fromstring(data: bytes):
    try:
        from defusedxml import ElementTree as SafeET  # noqa: PLC0415

        return SafeET.fromstring(data)
    except ModuleNotFoundError:
        import xml.etree.ElementTree as ET  # noqa: PLC0415

        return ET.fromstring(data)


def which(executable: str) -> str | None:
    return shutil.which(executable)


def require_input(path_text: str) -> Path:
    path = Path(path_text)
    if not path.is_file():
        raise ToolError(f"input not found or not a file: {path}", EXIT_INPUT)
    if path.suffix.lower() not in PPTX_EXTENSIONS:
        raise ToolError(
            f"unsupported presentation format {path.suffix or '(none)'}; use PPTX, or a verified read-only POTX/PPSX route",
            EXIT_UNSUPPORTED,
        )
    validate_archive_safety(path)
    return path


def require_regular_file(path_text: str, label: str) -> Path:
    path = Path(path_text)
    if not path.is_file():
        raise ToolError(f"{label} not found or not a file: {path}", EXIT_INPUT)
    return path


def safe_member_name(name: str) -> bool:
    normalized = name.replace("\\", "/")
    return (
        bool(normalized)
        and not normalized.startswith("/")
        and not re.match(r"^[A-Za-z]:", normalized)
        and ".." not in normalized.split("/")
    )


def validate_archive_safety(path: Path) -> None:
    if path.stat().st_size > MAX_ARCHIVE_BYTES:
        raise ToolError("presentation exceeds the 64 MiB compressed-size safety limit", EXIT_INPUT)
    try:
        archive = zipfile.ZipFile(path)
    except (zipfile.BadZipFile, OSError) as exc:
        raise ToolError(f"not a readable PPTX/ZIP package: {exc}", EXIT_INPUT) from exc
    with archive:
        infos = archive.infolist()
        if len(infos) > MAX_ENTRIES:
            raise ToolError("presentation contains too many ZIP members", EXIT_INPUT)
        names: set[str] = set()
        expanded = 0
        for info in infos:
            if not safe_member_name(info.filename):
                raise ToolError(f"unsafe ZIP member path: {info.filename}", EXIT_INPUT)
            if info.flag_bits & 0x1:
                raise ToolError(f"encrypted ZIP member is unsupported: {info.filename}", EXIT_UNSUPPORTED)
            mode = (info.external_attr >> 16) & 0o170000
            if mode == 0o120000:
                raise ToolError(f"symlink ZIP member is unsupported: {info.filename}", EXIT_INPUT)
            if info.filename in names:
                raise ToolError(f"duplicate ZIP member is unsupported: {info.filename}", EXIT_INPUT)
            names.add(info.filename)
            expanded += info.file_size
            if expanded > MAX_EXPANDED_BYTES:
                raise ToolError("expanded presentation exceeds the 256 MiB safety limit", EXIT_INPUT)
            if info.compress_size and info.file_size / info.compress_size > MAX_COMPRESSION_RATIO:
                raise ToolError(f"suspicious compression ratio in ZIP member: {info.filename}", EXIT_INPUT)
            if info.filename.endswith((".xml", ".rels")):
                if info.file_size > MAX_XML_BYTES:
                    raise ToolError(f"XML part exceeds the 16 MiB safety limit: {info.filename}", EXIT_INPUT)
                xml = archive.read(info)
                if re.search(br"<!DOCTYPE|<!ENTITY", xml, flags=re.I):
                    raise ToolError(
                        f"DTD or entity declaration is unsupported in package XML: {info.filename}",
                        EXIT_INPUT,
                    )
                if re.search(br"<(?:[A-Za-z0-9_]+:)?modifyVerifier\b", xml, flags=re.I):
                    raise ToolError("protected presentation packages are unsupported", EXIT_UNSUPPORTED)
        lower_names = {name.lower() for name in names}
        if "ppt/vbaproject.bin" in lower_names:
            raise ToolError("macro-bearing presentations are unsupported", EXIT_UNSUPPORTED)
        if any(name.startswith("_xmlsignatures/") for name in lower_names):
            raise ToolError("signed presentation packages are unsupported", EXIT_UNSUPPORTED)
        corrupt_member = archive.testzip()
        if corrupt_member:
            raise ToolError(f"ZIP integrity check failed at member: {corrupt_member}", EXIT_INPUT)


def package_feature_flags(path: Path) -> dict[str, int]:
    external_relationships = 0
    embedded_objects = 0
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        embedded_objects = sum(name.startswith("ppt/embeddings/") and not name.endswith("/") for name in names)
        for name in names:
            if not name.endswith(".rels"):
                continue
            root = xml_fromstring(archive.read(name))
            external_relationships += sum(
                relationship.get("TargetMode", "Internal") == "External"
                for relationship in root
            )
    return {
        "external_relationships": external_relationships,
        "embedded_objects": embedded_objects,
    }


def require_safe_renderer_input(path: Path) -> None:
    flags = package_feature_flags(path)
    active = [name for name, count in flags.items() if count]
    if active:
        raise ToolError(
            "rendering is blocked for packages with " + ", ".join(active),
            EXIT_UNSUPPORTED,
        )


def require_mutable_pptx(path: Path) -> None:
    if path.suffix.lower() != ".pptx":
        raise ToolError("mutating operations require a .pptx input and a distinct .pptx output", EXIT_UNSUPPORTED)


def require_distinct_output(input_path: Path, output_path: Path) -> None:
    try:
        same = input_path.resolve() == output_path.resolve()
    except OSError:
        same = input_path.absolute() == output_path.absolute()
    if same:
        raise ToolError("input and output paths must be distinct", EXIT_OUTPUT)


def check_output_path(path_text: str, overwrite: bool) -> Path:
    path = Path(path_text)
    if path.exists() and not overwrite:
        raise ToolError(
            f"output already exists: {path}; pass --overwrite only with explicit replacement approval",
            EXIT_OUTPUT,
        )
    if path.exists() and path.is_dir():
        raise ToolError(f"output path is a directory: {path}", EXIT_OUTPUT)
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def check_output_dir(path_text: str) -> Path:
    path = Path(path_text)
    if path.exists():
        if not path.is_dir():
            raise ToolError(f"output directory path is a file: {path}", EXIT_OUTPUT)
        if any(path.iterdir()):
            raise ToolError(
                f"output directory is not empty: {path}; choose a fresh directory so stale files cannot pose as results",
                EXIT_OUTPUT,
            )
    else:
        path.mkdir(parents=True)
    return path


def atomic_write_bytes(target: Path, data: bytes) -> None:
    handle = tempfile.NamedTemporaryFile(
        dir=str(target.parent), prefix=f".{target.name}.", suffix=".tmp", delete=False
    )
    try:
        handle.write(data)
        handle.close()
        os.replace(handle.name, target)
    except BaseException:
        handle.close()
        Path(handle.name).unlink(missing_ok=True)
        raise


def atomic_write_text(target: Path, text: str) -> None:
    atomic_write_bytes(target, text.encode("utf-8"))


def parse_slide_spec(spec: str, total: int) -> list[int]:
    """Parse a 1-based slide specification preserving order.

    Grammar: all | odd | even | N | N-M | last, comma-separated.
    """
    spec = (spec or "").strip().lower()
    if not spec:
        raise ToolError("empty slide specification", EXIT_ARGS)
    if spec == "all":
        return list(range(1, total + 1))
    if spec == "odd":
        return list(range(1, total + 1, 2))
    if spec == "even":
        return list(range(2, total + 1, 2))

    def resolve(token: str) -> int:
        if token == "last":
            return total
        if not token.isdigit():
            raise ToolError(f"invalid slide token: {token!r}", EXIT_ARGS)
        value = int(token)
        if value < 1:
            raise ToolError(f"slide numbers are 1-based; got {value}", EXIT_ARGS)
        if value > total:
            raise ToolError(f"slide {value} beyond deck of {total} slides", EXIT_ARGS)
        return value

    selected: list[int] = []
    for chunk in spec.split(","):
        chunk = chunk.strip()
        if not chunk:
            raise ToolError("empty token in slide specification", EXIT_ARGS)
        if "-" in chunk and chunk != "last":
            start_text, _, end_text = chunk.partition("-")
            start, end = resolve(start_text.strip()), resolve(end_text.strip())
            if end < start:
                raise ToolError(f"descending range not allowed: {chunk}", EXIT_ARGS)
            selected.extend(range(start, end + 1))
        else:
            selected.append(resolve(chunk))
    return selected


def result_envelope(command: str, backend: list[str], inputs: list[str]) -> dict:
    return {
        "ok": True,
        "command": command,
        "runtime": "python",
        "backend": backend,
        "inputs": inputs,
        "outputs": [],
        "warnings": [],
        "details": {},
    }


def emit(result: dict, args: argparse.Namespace) -> None:
    if args.json:
        print(json.dumps(result, indent=2, ensure_ascii=False))
    elif not args.quiet:
        for warning in result["warnings"]:
            print(f"warning: {warning}", file=sys.stderr)
        for output in result["outputs"]:
            print(output)


# ---------------------------------------------------------------------------
# Shape/text traversal helpers
# ---------------------------------------------------------------------------


def iter_shape_lines(shapes):
    """Yield text lines from shapes in package order, recursing into groups and tables."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shape_lines(shape.shapes)
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                yield "| " + " | ".join(cells) + " |"
            continue
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    prefix = "  " * min(paragraph.level, 4)
                    yield f"{prefix}{text}"


def iter_text_frames(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_frames(shape.shapes)
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
            continue
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            yield shape.text_frame


def notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def slide_title(slide) -> str:
    try:
        title_shape = slide.shapes.title
    except Exception:  # noqa: BLE001 - malformed layouts must not abort inspection
        title_shape = None
    if title_shape is not None and title_shape.has_text_frame:
        return title_shape.text_frame.text.strip()
    return ""


# ---------------------------------------------------------------------------
# Native tool wrappers (LibreOffice + Poppler)
# ---------------------------------------------------------------------------


def run_process(argv: list[str], timeout: int = 300) -> subprocess.CompletedProcess:
    try:
        return subprocess.run(
            argv, capture_output=True, text=True, timeout=timeout, check=False
        )
    except FileNotFoundError as exc:
        raise ToolError(f"required executable not found: {argv[0]}", EXIT_MISSING_DEPENDENCY) from exc
    except subprocess.TimeoutExpired as exc:
        raise ToolError(f"{argv[0]} timed out after {timeout}s", EXIT_OUTPUT) from exc


def soffice_convert_to_pdf(input_path: Path, out_dir: Path) -> Path:
    if which("soffice") is None:
        raise ToolError("LibreOffice (soffice) is required for rendering", EXIT_MISSING_DEPENDENCY)
    with tempfile.TemporaryDirectory(prefix="pptx-soffice-profile-") as profile:
        completed = run_process(
            [
                "soffice",
                "--headless",
                "--norestore",
                f"-env:UserInstallation=file://{profile}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(input_path),
            ],
            timeout=600,
        )
    pdf_path = out_dir / (input_path.stem + ".pdf")
    if completed.returncode != 0 or not pdf_path.is_file():
        detail = (completed.stderr or completed.stdout or "").strip()[-500:]
        raise ToolError(f"LibreOffice conversion failed: {detail or 'no diagnostic output'}", EXIT_OUTPUT)
    return pdf_path


def pdftoppm_render(pdf_path: Path, out_dir: Path, dpi: int, image_format: str) -> list[Path]:
    if which("pdftoppm") is None:
        raise ToolError("Poppler (pdftoppm) is required for rendering", EXIT_MISSING_DEPENDENCY)
    flag = "-png" if image_format == "png" else "-jpeg"
    completed = run_process(
        ["pdftoppm", flag, "-r", str(dpi), str(pdf_path), str(out_dir / "slide")],
        timeout=600,
    )
    if completed.returncode != 0:
        raise ToolError(f"pdftoppm failed: {(completed.stderr or '').strip()[-500:]}", EXIT_OUTPUT)
    suffix = ".png" if image_format == "png" else ".jpg"
    images = sorted(out_dir.glob(f"slide-*{suffix}"))
    if not images:
        raise ToolError("rendering produced no page images", EXIT_OUTPUT)
    return images


def render_slides_to_images(
    input_path: Path, out_dir: Path, dpi: int, image_format: str, spec: str | None
) -> tuple[list[Path], list[int]]:
    with tempfile.TemporaryDirectory(prefix="pptx-render-") as scratch_text:
        scratch = Path(scratch_text)
        pdf_path = soffice_convert_to_pdf(input_path, scratch)
        images = pdftoppm_render(pdf_path, scratch, dpi, image_format)
        total = len(images)
        selected = parse_slide_spec(spec, total) if spec else list(range(1, total + 1))
        if len(set(selected)) != len(selected):
            raise ToolError("repeated slides are not allowed in a render specification", EXIT_ARGS)
        width = max(2, len(str(total)))
        kept: list[Path] = []
        suffix = ".png" if image_format == "png" else ".jpg"
        for slide_number in selected:
            source = images[slide_number - 1]
            target = out_dir / f"slide-{slide_number:0{width}d}{suffix}"
            shutil.move(str(source), target)
            kept.append(target)
        return kept, selected


def write_directory_manifest(out_dir: Path, command: str, entries: list[dict], ok: bool, error: str | None = None) -> None:
    manifest = {"ok": ok, "command": command, "runtime": "python", "outputs": entries}
    if error:
        manifest["error"] = error
    (out_dir / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------


def cmd_doctor(args: argparse.Namespace) -> dict:
    def module_present(name: str) -> bool:
        import importlib.util  # noqa: PLC0415

        return importlib.util.find_spec(name) is not None

    modules = {name: module_present(name) for name in ("pptx", "PIL", "defusedxml")}
    native = {name: which(name) is not None for name in ("soffice", "pdftoppm")}
    commands = {
        "inspect": modules["pptx"],
        "extract-text": modules["pptx"],
        "extract-notes": modules["pptx"],
        "extract-media": True,
        "select": modules["pptx"],
        "replace-text": modules["pptx"],
        "check": True,
        "render": native["soffice"] and native["pdftoppm"],
        "contact-sheet": native["soffice"] and native["pdftoppm"] and modules["PIL"],
    }
    result = result_envelope("doctor", ["python"], [])
    result["details"] = {
        "python": sys.version.split()[0],
        "modules": modules,
        "native_tools": native,
        "command_readiness": commands,
    }
    if not args.json and not args.quiet:
        print(f"python {result['details']['python']}")
        for name, present in {**modules, **native}.items():
            print(f"{'ok  ' if present else 'MISS'} {name}")
        for name, ready in commands.items():
            print(f"{'ready  ' if ready else 'blocked'} {name}")
    return result


def cmd_inspect(args: argparse.Namespace) -> dict:
    pptx = load_pptx_module()
    input_path = require_input(args.input)
    try:
        presentation = pptx.Presentation(str(input_path))
    except Exception as exc:  # noqa: BLE001 - report unreadable package
        raise ToolError(f"cannot open as a PPTX package: {exc}", EXIT_INPUT) from exc

    emu_per_inch = 914400
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        slides.append(
            {
                "number": index,
                "layout": slide.slide_layout.name,
                "title": slide_title(slide)[:120],
                "shape_count": len(slide.shapes),
                "has_notes": bool(notes_text(slide)),
            }
        )
    package = presentation.part.package
    media_parts = [
        str(part.partname)
        for part in package.iter_parts()
        if str(part.partname).startswith("/ppt/media/")
    ]
    core = presentation.core_properties
    feature_flags = package_feature_flags(input_path)
    result = result_envelope("inspect", ["python-pptx"], [str(input_path)])
    result["details"] = {
        "slide_count": len(presentation.slides._sldIdLst),
        "slide_width_in": round(presentation.slide_width / emu_per_inch, 3),
        "slide_height_in": round(presentation.slide_height / emu_per_inch, 3),
        "layout_count": sum(len(master.slide_layouts) for master in presentation.slide_masters),
        "master_count": len(presentation.slide_masters),
        "media_count": len(media_parts),
        "core_properties": {
            "title": core.title or "",
            "author": core.author or "",
            "last_modified_by": core.last_modified_by or "",
            "modified": core.modified.isoformat() if core.modified else "",
        },
        "package_flags": feature_flags,
        "slides": slides,
    }
    if not args.json and not args.quiet:
        details = result["details"]
        print(
            f"slides={details['slide_count']} size={details['slide_width_in']}x{details['slide_height_in']}in "
            f"masters={details['master_count']} layouts={details['layout_count']} media={details['media_count']}"
        )
        for slide in slides:
            marker = " +notes" if slide["has_notes"] else ""
            print(f"  {slide['number']:>3} [{slide['layout']}] {slide['title']}{marker}")
    return result


def cmd_extract_text(args: argparse.Namespace) -> dict:
    pptx = load_pptx_module()
    input_path = require_input(args.input)
    output_path = Path(args.output)
    require_distinct_output(input_path, output_path)
    output_path = check_output_path(args.output, args.overwrite)
    presentation = pptx.Presentation(str(input_path))
    slides = list(presentation.slides)
    selected = parse_slide_spec(args.slides, len(slides)) if args.slides else list(range(1, len(slides) + 1))

    blocks: list[str] = []
    for number in selected:
        slide = slides[number - 1]
        lines = list(iter_shape_lines(slide.shapes))
        block = [f"## Slide {number} — {slide.slide_layout.name}", ""]
        block.extend(lines or ["(no text)"])
        if args.with_notes:
            notes = notes_text(slide)
            if notes:
                block.extend(["", "> Notes: " + notes.replace("\n", "\n> ")])
        blocks.append("\n".join(block))
    atomic_write_text(output_path, "\n\n".join(blocks) + "\n")

    result = result_envelope("extract-text", ["python-pptx"], [str(input_path)])
    result["outputs"] = [str(output_path)]
    result["details"] = {"slides_extracted": selected}
    return result


def cmd_extract_notes(args: argparse.Namespace) -> dict:
    pptx = load_pptx_module()
    input_path = require_input(args.input)
    output_path = Path(args.output)
    require_distinct_output(input_path, output_path)
    output_path = check_output_path(args.output, args.overwrite)
    presentation = pptx.Presentation(str(input_path))
    slides = list(presentation.slides)
    selected = parse_slide_spec(args.slides, len(slides)) if args.slides else list(range(1, len(slides) + 1))

    blocks: list[str] = []
    slides_with_notes = 0
    for number in selected:
        notes = notes_text(slides[number - 1])
        if notes:
            slides_with_notes += 1
        blocks.append(f"## Slide {number}\n\n{notes or '(no notes)'}")
    atomic_write_text(output_path, "\n\n".join(blocks) + "\n")

    result = result_envelope("extract-notes", ["python-pptx"], [str(input_path)])
    result["outputs"] = [str(output_path)]
    result["details"] = {"slides_checked": selected, "slides_with_notes": slides_with_notes}
    return result


def cmd_extract_media(args: argparse.Namespace) -> dict:
    input_path = require_input(args.input)
    out_dir = check_output_dir(args.output_dir)
    entries: list[dict] = []
    try:
        with zipfile.ZipFile(input_path) as archive:
            media_names = [
                name for name in archive.namelist() if name.startswith("ppt/media/") and not name.endswith("/")
            ]
            for name in media_names:
                data = archive.read(name)
                target = out_dir / Path(name).name
                atomic_write_bytes(target, data)
                entries.append({"path": str(target), "bytes": len(data), "source_part": "/" + name})
    except zipfile.BadZipFile as exc:
        raise ToolError(f"not a readable PPTX/ZIP package: {exc}", EXIT_INPUT) from exc

    write_directory_manifest(out_dir, "extract-media", entries, ok=True)
    result = result_envelope("extract-media", ["zipfile"], [str(input_path)])
    result["outputs"] = [entry["path"] for entry in entries]
    result["details"] = {"media_count": len(entries), "manifest": str(out_dir / "manifest.json")}
    if not entries:
        result["warnings"].append("deck contains no media parts")
    return result


def cmd_select(args: argparse.Namespace) -> dict:
    pptx = load_pptx_module()
    input_path = require_input(args.input)
    require_mutable_pptx(input_path)
    output_path = Path(args.output)
    require_distinct_output(input_path, output_path)
    output_path = check_output_path(args.output, args.overwrite)
    if output_path.suffix.lower() != ".pptx":
        raise ToolError("select output must use the .pptx extension", EXIT_ARGS)
    presentation = pptx.Presentation(str(input_path))
    sld_id_lst = presentation.slides._sldIdLst
    slide_ids = list(sld_id_lst)
    total = len(slide_ids)
    selected = parse_slide_spec(args.slides, total)
    if not selected:
        raise ToolError("selection must keep at least one slide", EXIT_ARGS)
    if len(set(selected)) != len(selected):
        raise ToolError(
            "slide duplication is not supported by select; repeated slide numbers were rejected",
            EXIT_UNSUPPORTED,
        )

    keep = set(selected)
    for position, slide_id in enumerate(slide_ids, start=1):
        if position not in keep:
            rel_id = slide_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            sld_id_lst.remove(slide_id)
            presentation.part.drop_rel(rel_id)
    # Reorder the surviving entries to the requested order.
    survivors = {position: slide_id for position, slide_id in enumerate(slide_ids, start=1) if position in keep}
    for position in selected:
        sld_id_lst.append(survivors[position])  # lxml append moves the element

    with tempfile.NamedTemporaryFile(
        dir=str(output_path.parent), prefix=f".{output_path.name}.", suffix=".tmp", delete=False
    ) as handle:
        temp_name = handle.name
    try:
        presentation.save(temp_name)
        reopened = pptx.Presentation(temp_name)
        if len(reopened.slides._sldIdLst) != len(selected):
            raise ToolError("saved deck does not contain the requested slide count", EXIT_OUTPUT)
        os.replace(temp_name, output_path)
    except BaseException:
        Path(temp_name).unlink(missing_ok=True)
        raise

    result = result_envelope("select", ["python-pptx"], [str(input_path)])
    result["outputs"] = [str(output_path)]
    result["details"] = {"kept_order": selected, "source_slide_count": total}
    return result


def load_replacement_map(path: Path) -> list[tuple[str, str]]:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise ToolError(f"replacement map is not valid JSON: {exc}", EXIT_ARGS) from exc
    replacements = data.get("replacements", data) if isinstance(data, dict) else data
    pairs: list[tuple[str, str]] = []
    if isinstance(replacements, dict):
        pairs = [(str(find), str(replace)) for find, replace in replacements.items()]
    elif isinstance(replacements, list):
        for item in replacements:
            if not isinstance(item, dict) or "find" not in item or "replace" not in item:
                raise ToolError("each replacement entry needs 'find' and 'replace'", EXIT_ARGS)
            pairs.append((str(item["find"]), str(item["replace"])))
    if not pairs or any(not find for find, _ in pairs):
        raise ToolError("replacement map must contain at least one non-empty 'find' string", EXIT_ARGS)
    return pairs


def cmd_replace_text(args: argparse.Namespace) -> dict:
    pptx = load_pptx_module()
    input_path = require_input(args.input)
    require_mutable_pptx(input_path)
    output_path = Path(args.output)
    require_distinct_output(input_path, output_path)
    map_path = require_regular_file(args.map, "replacement map")
    require_distinct_output(map_path, output_path)
    output_path = check_output_path(args.output, args.overwrite)
    if output_path.suffix.lower() != ".pptx":
        raise ToolError("replace-text output must use the .pptx extension", EXIT_ARGS)
    pairs = load_replacement_map(map_path)
    presentation = pptx.Presentation(str(input_path))
    slides = list(presentation.slides)
    selected = parse_slide_spec(args.slides, len(slides)) if args.slides else list(range(1, len(slides) + 1))

    counts = {find: 0 for find, _ in pairs}
    fragmented: set[str] = set()
    for number in selected:
        slide = slides[number - 1]
        frames = list(iter_text_frames(slide.shapes))
        if args.include_notes and slide.has_notes_slide:
            frames.append(slide.notes_slide.notes_text_frame)
        for frame in frames:
            for paragraph in frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs)
                for find, replace in pairs:
                    hits_here = 0
                    for run in paragraph.runs:
                        if find in run.text:
                            hits_here += run.text.count(find)
                            run.text = run.text.replace(find, replace)
                    counts[find] += hits_here
                    if hits_here == 0 and find in paragraph_text:
                        fragmented.add(find)

    with tempfile.NamedTemporaryFile(
        dir=str(output_path.parent), prefix=f".{output_path.name}.", suffix=".tmp", delete=False
    ) as handle:
        temp_name = handle.name
    try:
        presentation.save(temp_name)
        os.replace(temp_name, output_path)
    except BaseException:
        Path(temp_name).unlink(missing_ok=True)
        raise

    result = result_envelope("replace-text", ["python-pptx"], [str(input_path)])
    result["outputs"] = [str(output_path)]
    result["details"] = {"replacements": counts, "slides_processed": selected}
    for find, count in counts.items():
        if count == 0:
            result["warnings"].append(f"no replacement made for {find!r}")
    for find in sorted(fragmented):
        result["warnings"].append(
            f"{find!r} spans multiple runs in at least one paragraph; run-level replacement could not rewrite it"
        )
    return result


def cmd_render(args: argparse.Namespace) -> dict:
    image_format = args.format
    if not 36 <= args.dpi <= 600:
        raise ToolError("--dpi must be from 36 to 600", EXIT_ARGS)
    input_path = require_input(args.input)
    require_safe_renderer_input(input_path)
    out_dir = check_output_dir(args.output_dir)
    try:
        images, selected = render_slides_to_images(input_path, out_dir, args.dpi, image_format, args.slides)
    except ToolError as exc:
        write_directory_manifest(out_dir, "render", [], ok=False, error=str(exc))
        raise
    entries = [
        {"path": str(path), "slide": slide_number}
        for path, slide_number in zip(images, selected)
    ]
    write_directory_manifest(out_dir, "render", entries, ok=True)
    result = result_envelope("render", ["soffice", "pdftoppm"], [str(input_path)])
    result["outputs"] = [entry["path"] for entry in entries]
    result["details"] = {
        "dpi": args.dpi,
        "format": image_format,
        "slides_rendered": selected,
        "manifest": str(out_dir / "manifest.json"),
    }
    return result


def cmd_contact_sheet(args: argparse.Namespace) -> dict:
    columns = args.columns
    if columns < 1:
        raise ToolError("--columns must be at least 1", EXIT_ARGS)
    if columns > 20:
        raise ToolError("--columns must not exceed 20", EXIT_ARGS)
    if not 64 <= args.width <= 4000:
        raise ToolError("--width must be from 64 to 4000 pixels", EXIT_ARGS)
    if not 36 <= args.dpi <= 600:
        raise ToolError("--dpi must be from 36 to 600", EXIT_ARGS)
    Image, ImageDraw = load_pillow_module()
    input_path = require_input(args.input)
    require_safe_renderer_input(input_path)
    output_path = Path(args.output)
    require_distinct_output(input_path, output_path)
    output_path = check_output_path(args.output, args.overwrite)

    with tempfile.TemporaryDirectory(prefix="pptx-sheet-") as scratch_text:
        scratch = Path(scratch_text)
        images, selected = render_slides_to_images(input_path, scratch, args.dpi, "png", None)
        thumbs = []
        for path in images:
            image = Image.open(path).convert("RGB")
            ratio = args.width / image.width
            thumbs.append(image.resize((args.width, max(1, int(image.height * ratio)))))
        rows = (len(thumbs) + columns - 1) // columns
        pad, label_height = 10, 22
        cell_height = max(thumb.height for thumb in thumbs) + label_height
        sheet = Image.new(
            "RGB",
            (columns * (args.width + pad) + pad, rows * (cell_height + pad) + pad),
            (245, 245, 245),
        )
        draw = ImageDraw.Draw(sheet)
        for index, thumb in enumerate(thumbs):
            column, row = index % columns, index // columns
            x = pad + column * (args.width + pad)
            y = pad + row * (cell_height + pad)
            sheet.paste(thumb, (x, y))
            draw.text((x + 2, y + thumb.height + 4), f"Slide {selected[index]}", fill=(20, 20, 20))
        with tempfile.NamedTemporaryFile(
            dir=str(output_path.parent), prefix=f".{output_path.name}.", suffix=".tmp", delete=False
        ) as handle:
            temp_name = handle.name
        try:
            sheet.save(temp_name, format="PNG")
            os.replace(temp_name, output_path)
        except BaseException:
            Path(temp_name).unlink(missing_ok=True)
            raise

    result = result_envelope("contact-sheet", ["soffice", "pdftoppm", "Pillow"], [str(input_path)])
    result["outputs"] = [str(output_path)]
    result["details"] = {"slide_count": len(selected), "columns": columns, "thumb_width": args.width}
    return result


CT_NAMESPACE = "{http://schemas.openxmlformats.org/package/2006/content-types}"
REL_NAMESPACE = "{http://schemas.openxmlformats.org/package/2006/relationships}"
PRESENTATION_NAMESPACE = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
OFFICE_REL_NAMESPACE = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def cmd_check(args: argparse.Namespace) -> dict:
    input_path = require_input(args.input)
    errors: list[str] = []
    warnings: list[str] = []

    try:
        archive = zipfile.ZipFile(input_path)
    except zipfile.BadZipFile as exc:
        raise ToolError(f"not a readable ZIP package: {exc}", EXIT_INPUT) from exc

    with archive:
        names = set(archive.namelist())
        corrupt_member = archive.testzip()
        if corrupt_member:
            errors.append(f"corrupt ZIP member: {corrupt_member}")

        # Content types must parse and cover every part.
        defaults: dict[str, str] = {}
        overrides: dict[str, str] = {}
        if "[Content_Types].xml" not in names:
            errors.append("missing [Content_Types].xml")
        else:
            try:
                root = xml_fromstring(archive.read("[Content_Types].xml"))
                for node in root:
                    if node.tag == f"{CT_NAMESPACE}Default":
                        defaults[node.get("Extension", "").lower()] = node.get("ContentType", "")
                    elif node.tag == f"{CT_NAMESPACE}Override":
                        overrides[node.get("PartName", "")] = node.get("ContentType", "")
            except Exception as exc:  # noqa: BLE001
                errors.append(f"[Content_Types].xml does not parse: {exc}")

        for name in sorted(names):
            if name.endswith("/") or name == "[Content_Types].xml":
                continue
            extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
            if f"/{name}" not in overrides and extension not in defaults:
                errors.append(f"part has no declared content type: /{name}")

        # Every XML part must be well-formed; every internal rel target must exist.
        rels_files = [name for name in names if name.endswith(".rels")]
        for name in sorted(names):
            if not (name.endswith(".xml") or name.endswith(".rels")):
                continue
            try:
                xml_fromstring(archive.read(name))
            except Exception as exc:  # noqa: BLE001
                errors.append(f"XML part is not well-formed: /{name}: {exc}")

        for rels_name in sorted(rels_files):
            base = posixpath.dirname(posixpath.dirname(rels_name))
            try:
                root = xml_fromstring(archive.read(rels_name))
            except Exception:  # noqa: BLE001 - already reported above
                continue
            for relationship in root:
                if relationship.get("TargetMode", "Internal") == "External":
                    continue
                target = relationship.get("Target", "")
                resolved = posixpath.normpath(posixpath.join(base, target)).lstrip("/")
                if resolved not in names:
                    errors.append(f"{rels_name}: relationship target missing: {target}")

        # The slide list must resolve to real slide parts.
        slide_count = 0
        if "ppt/presentation.xml" in names and "ppt/_rels/presentation.xml.rels" in names:
            try:
                presentation_root = xml_fromstring(archive.read("ppt/presentation.xml"))
                rels_root = xml_fromstring(archive.read("ppt/_rels/presentation.xml.rels"))
                rel_targets = {
                    rel.get("Id"): rel.get("Target", "") for rel in rels_root
                }
                sld_id_lst = presentation_root.find(f"{PRESENTATION_NAMESPACE}sldIdLst")
                for sld_id in sld_id_lst if sld_id_lst is not None else []:
                    slide_count += 1
                    rel_id = sld_id.get(f"{OFFICE_REL_NAMESPACE}id")
                    target = rel_targets.get(rel_id)
                    if target is None:
                        errors.append(f"sldIdLst references undeclared relationship {rel_id}")
                        continue
                    resolved = posixpath.normpath(posixpath.join("ppt", target)).lstrip("/")
                    if resolved not in names:
                        errors.append(f"slide part missing for {rel_id}: {target}")
            except Exception as exc:  # noqa: BLE001
                errors.append(f"presentation.xml structure check failed: {exc}")
        else:
            errors.append("missing ppt/presentation.xml or its relationships part")

        if slide_count == 0 and not errors:
            warnings.append("deck has no slides in sldIdLst")

    # Confirm the modeled library can also open the package when it is available.
    try:
        pptx = load_pptx_module()
        pptx.Presentation(str(input_path))
    except ToolError:
        warnings.append("python-pptx unavailable; library open check skipped")
    except Exception as exc:  # noqa: BLE001
        errors.append(f"python-pptx cannot open the package: {exc}")

    feature_flags = package_feature_flags(input_path)
    if feature_flags["external_relationships"]:
        warnings.append(
            f"package contains {feature_flags['external_relationships']} external relationship(s); they were not followed"
        )
    if feature_flags["embedded_objects"]:
        warnings.append(
            f"package contains {feature_flags['embedded_objects']} embedded object(s); they were not executed"
        )
    result = result_envelope("check", ["zipfile", "xml"], [str(input_path)])
    result["warnings"] = warnings
    result["details"] = {"errors": errors, "slide_count": slide_count, "package_flags": feature_flags}
    if errors:
        result["ok"] = False
        if not args.json:
            for error in errors:
                print(f"error: {error}", file=sys.stderr)
        raise ToolCheckFailure(result)
    if not args.json and not args.quiet:
        print(f"check passed: {slide_count} slides, no structural errors")
    return result


class ToolCheckFailure(Exception):
    def __init__(self, result: dict) -> None:
        super().__init__("structural check failed")
        self.result = result


# ---------------------------------------------------------------------------
# CLI wiring
# ---------------------------------------------------------------------------


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="pptx-tool",
        description="Python PPTX backend for q-tool-pptx",
        epilog="Slide numbers are 1-based. Outputs refuse existing targets without --overwrite.",
    )
    parser.add_argument("--json", action="store_true", help="write one JSON result object to stdout")
    parser.add_argument("--quiet", action="store_true", help="suppress non-essential diagnostics")
    parser.add_argument("--overwrite", action="store_true", help="explicit approval to replace an existing output file")
    subparsers = parser.add_subparsers(dest="command", metavar="COMMAND")

    sub = subparsers.add_parser("doctor")

    sub = subparsers.add_parser("inspect")
    sub.add_argument("input")

    sub = subparsers.add_parser("extract-text")
    sub.add_argument("input")
    sub.add_argument("--output", required=True)
    sub.add_argument("--slides")
    sub.add_argument("--with-notes", action="store_true")

    sub = subparsers.add_parser("extract-notes")
    sub.add_argument("input")
    sub.add_argument("--output", required=True)
    sub.add_argument("--slides")

    sub = subparsers.add_parser("extract-media")
    sub.add_argument("input")
    sub.add_argument("--output-dir", required=True)

    sub = subparsers.add_parser("select")
    sub.add_argument("input")
    sub.add_argument("--slides", required=True)
    sub.add_argument("--output", required=True)

    sub = subparsers.add_parser("replace-text")
    sub.add_argument("input")
    sub.add_argument("--map", required=True)
    sub.add_argument("--output", required=True)
    sub.add_argument("--slides")
    sub.add_argument("--include-notes", action="store_true")

    sub = subparsers.add_parser("render")
    sub.add_argument("input")
    sub.add_argument("--output-dir", required=True)
    sub.add_argument("--dpi", type=int, default=150)
    sub.add_argument("--format", choices=("png", "jpeg"), default="png")
    sub.add_argument("--slides")

    sub = subparsers.add_parser("contact-sheet")
    sub.add_argument("input")
    sub.add_argument("--output", required=True)
    sub.add_argument("--columns", type=int, default=4)
    sub.add_argument("--width", type=int, default=320)
    sub.add_argument("--dpi", type=int, default=80)

    sub = subparsers.add_parser("check")
    sub.add_argument("input")

    return parser


COMMAND_HANDLERS = {
    "doctor": cmd_doctor,
    "inspect": cmd_inspect,
    "extract-text": cmd_extract_text,
    "extract-notes": cmd_extract_notes,
    "extract-media": cmd_extract_media,
    "select": cmd_select,
    "replace-text": cmd_replace_text,
    "render": cmd_render,
    "contact-sheet": cmd_contact_sheet,
    "check": cmd_check,
}


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    if not args.command:
        parser.print_help()
        return EXIT_OK
    handler = COMMAND_HANDLERS[args.command]
    try:
        result = handler(args)
    except ToolCheckFailure as failure:
        if args.json:
            print(json.dumps(failure.result, indent=2, ensure_ascii=False))
        return EXIT_OUTPUT
    except ToolError as exc:
        if args.json:
            print(
                json.dumps(
                    {"ok": False, "command": args.command, "runtime": "python", "error": str(exc)},
                    ensure_ascii=False,
                )
            )
        print(f"Error: {exc}", file=sys.stderr)
        return exc.code
    emit(result, args)
    return EXIT_OK


if __name__ == "__main__":
    raise SystemExit(main())
