#!/usr/bin/env python3
"""Build deterministic PPTX fixtures for the q-tool-pptx smoke tests.

Creates deck.pptx (4 slides: title, bullets with notes, table, picture) and a
replacement map used by the replace-text smoke checks. Content is fixed so
assertions can match exact strings.
"""

from __future__ import annotations

import io
import json
import sys
import zipfile
from pathlib import Path


def build_deck(target: Path) -> None:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()

    # Slide 1: title + subtitle.
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Quasar Fixture Deck"
    slide.placeholders[1].text = "Deterministic fixture for smoke tests"

    # Slide 2: bullets + speaker notes.
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Agenda"
    body = slide.placeholders[1].text_frame
    body.text = "First bullet FIXTURE-TOKEN"
    for text in ("Second bullet", "Third bullet"):
        paragraph = body.add_paragraph()
        paragraph.text = text
        paragraph.level = 1
    fragmented = body.add_paragraph()
    fragmented.add_run().text = "{{FRAG"
    fragmented.add_run().text = "MENTED}}"
    slide.notes_slide.notes_text_frame.text = "Speaker notes for slide two"

    # Slide 3: a 2x3 table.
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Metrics"
    table = slide.shapes.add_table(2, 3, Inches(1), Inches(2), Inches(8), Inches(2)).table
    for column, header in enumerate(("Metric", "Value", "Unit")):
        table.cell(0, column).text = header
    for column, value in enumerate(("Latency", "42", "ms")):
        table.cell(1, column).text = value

    # Slide 4: generated picture.
    image = Image.new("RGB", (320, 200), (30, 90, 160))
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Picture"
    slide.shapes.add_picture(buffer, Inches(2), Inches(2), width=Inches(4))
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(40)

    presentation.save(str(target))


def main() -> int:
    out_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("fixtures")
    out_dir.mkdir(parents=True, exist_ok=True)
    build_deck(out_dir / "deck.pptx")
    (out_dir / "replacements.json").write_text(
        json.dumps(
            {
                "replacements": {
                    "FIXTURE-TOKEN": "REPLACED-TOKEN",
                    "{{FRAGMENTED}}": "SHOULD-NOT-REPLACE",
                    "not-in-deck": "unused",
                }
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    source = out_dir / "deck.pptx"
    with zipfile.ZipFile(source) as original:
        members = [(info, original.read(info)) for info in original.infolist()]
    with zipfile.ZipFile(out_dir / "unsafe-member.pptx", "w", compression=zipfile.ZIP_DEFLATED) as unsafe:
        for info, data in members:
            unsafe.writestr(info, data)
        unsafe.writestr("../escape.xml", "<unsafe/>")
    with zipfile.ZipFile(out_dir / "macro-bearing.pptx", "w", compression=zipfile.ZIP_DEFLATED) as macro:
        for info, data in members:
            macro.writestr(info, data)
        macro.writestr("ppt/vbaProject.bin", b"not-a-real-vba-project")
    with zipfile.ZipFile(out_dir / "protected.pptx", "w", compression=zipfile.ZIP_DEFLATED) as protected:
        for info, data in members:
            if info.filename == "ppt/presentation.xml":
                data = data.replace(
                    b"</p:presentation>",
                    b'<p:modifyVerifier cryptProviderType="rsaAES" cryptAlgorithmClass="hash"/></p:presentation>',
                )
            protected.writestr(info, data)
    with zipfile.ZipFile(out_dir / "external-relationship.pptx", "w", compression=zipfile.ZIP_DEFLATED) as external:
        for info, data in members:
            if info.filename == "ppt/_rels/presentation.xml.rels":
                data = data.replace(
                    b"</Relationships>",
                    b'<Relationship Id="rExternal" Type="urn:quasar:test" Target="https://example.invalid/deck" TargetMode="External"/></Relationships>',
                )
            external.writestr(info, data)
    print(str(out_dir / "deck.pptx"))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
